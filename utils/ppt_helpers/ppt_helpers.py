from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


def create_summary_slide(prs, title="summary", assignee="Assignee/Contributors: ?", deadline="Deadline: ?",
                         krs="KRs: ?", stakeholder="Stakeholder: ?",
                         accomplishments="", challenges="", status="Status: ?", start="Start: ?",
                         estimate="Estimate: ?", scope_changed="Scope Changed: ?",
                         progress_statuses_dates=None, gann_chart_path=None):
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    if progress_statuses_dates is None:
        progress_statuses_dates = [("Status 1", "Date 1"), ("Status 2", "Date 2"), ("Status 3", "Date 3")]

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(240, 240, 240)
    background.line.fill.background()

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(25, 94, 95)
    header.line.fill.background()

    summary_title = slide.shapes.add_textbox(Inches(0.5), Inches(-0.25), Inches(9), Inches(1))
    summary_text_frame = summary_title.text_frame
    summary_text = summary_text_frame.add_paragraph()
    summary_text.text = title
    summary_text.font.size = Pt(44)
    summary_text.font.bold = True
    summary_text.font.color.rgb = RGBColor(255, 255, 255)

    self.add_textbox_with_label(slide, "Assignee/Contributors: ", assignee, 0.5, 0.9)
    self.add_textbox_with_label(slide, "Stakeholder: ", stakeholder, 7.5, 0.9)
    self.add_textbox_with_label(slide, "KRs: ", krs, 7.5, 1.4)
    self.add_textbox_with_label(slide, "Deadline: ", deadline, 0.5, 1.4)

    self.add_accomplishments_and_challenges(slide, accomplishments, challenges)

    self.add_status_info(slide, start, status, estimate, scope_changed)


def create_squad_slide(prs, title="InStore Squad"):
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(25, 94, 95)
    background.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(4), Inches(3), Inches(5), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = title
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(44)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
    paragraph.alignment = PP_ALIGN.CENTER


def create_assignee_slide(prs, assignee_name):
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(25, 94, 95)
    background.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(4), Inches(3), Inches(5), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = f"{assignee_name}"
    title_paragraph = text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255)
    title_paragraph.alignment = PP_ALIGN.CENTER


def insert_gantt_chart_to_slide(prs, image_path):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), height=Inches(5.5))


def add_textbox_with_label(slide, label, variable, left, top):
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(4), Inches(0.5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    run1 = text.add_run()
    run1.text = label
    run1.font.size = Pt(20)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0, 0, 0)
    run2 = text.add_run()
    run2.text = f"{variable}"
    run2.font.size = Pt(20)
    run2.font.bold = False
    run2.font.color.rgb = RGBColor(0, 0, 0)


def add_accomplishments_and_challenges(slide, accomplishments, challenges):
    add_textbox_with_label(slide, "Accomplishments: ", accomplishments, 0.5, 2.5)
    add_textbox_with_label(slide, "Challenges: ", challenges, 7.5, 2.5)


def add_status_info(slide, start, status, estimate, scope_changed):
    add_centered_textbox(slide, Inches(0.25), Inches(6.5), Inches(2.3), Inches(0.5), Inches(6.3), "Start: ",
                         start)
    add_centered_textbox(slide, Inches(2.7), Inches(6.5), Inches(2.3), Inches(0.5), Inches(6.3), "Status: ",
                         status)
    add_centered_textbox(slide, Inches(5.2), Inches(6.5), Inches(2.3), Inches(0.5), Inches(6.3), "Estimate: ",
                         estimate)
    add_centered_textbox(slide, Inches(7.7), Inches(6.5), Inches(2.3), Inches(0.5), Inches(6.3),
                         "Scope Changed: ", scope_changed)


def add_centered_textbox(slide, left, top, width, height, top_text_box, label, variable, font_size=Pt(18)):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.line.color.rgb = RGBColor(0, 0, 0)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)

    text_box = slide.shapes.add_textbox(left, top_text_box, width, height)
    text_frame = text_box.text_frame
    text_frame.clear()

    text_frame.vertical_anchor = MSO_SHAPE.RECTANGLE
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run1 = p.add_run()
    run1.text = label
    run1.font.size = font_size
    run1.font.bold = True

    run2 = p.add_run()
    run2.text = variable
    run2.font.size = font_size
    run2.font.bold = False
