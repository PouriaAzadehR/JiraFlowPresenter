from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()


# Function to create a slide template
def create_summary_slide(prs, title="summary", assignee="Assignee/Contributors: ?", deadline="Deadline: ?",
                         krs="KRs: ?", stakeholder="Stakeholder: ?",
                         accomplishments="", challenges="", status="Status: ?", start="Start: ?",
                         estimate="Estimate: ?", scope_changed="Scope Changed: ?",
                         progress_statuses_dates=None,gann_chart_path=None):
    # Adjust the slide dimensions to make it wider (e.g., 16:9 aspect ratio)
    prs.slide_width = Inches(13.33)  # 13.33 inches wide (16:9 aspect ratio)
    prs.slide_height = Inches(7.5)  # 7.5 inches tall (16:9 aspect ratio)

    # Add a slide with a custom layout
    if progress_statuses_dates is None:
        progress_statuses_dates = [("Status 1", "Date 1"), ("Status 2", "Date 2"),
                                   ("Status 3", "Date 3")]
    slide_layout = prs.slide_layouts[5]  # Use a blank layout as a starting point
    slide = prs.slides.add_slide(slide_layout)

    # Add a background color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray background
    background.line.fill.background()  # No border

    # Add a header background for the title
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(25, 94, 95)  # Orange header background
    header.line.fill.background()  # No border

    # Add "summary" title
    summary_title_left = Inches(0.5)
    summary_title_top = Inches(-0.25)  # Move title to the top of the slide
    summary_title_width = Inches(9)
    summary_title_height = Inches(1)
    summary_title = slide.shapes.add_textbox(summary_title_left, summary_title_top, summary_title_width,
                                             summary_title_height)
    summary_text_frame = summary_title.text_frame
    summary_text = summary_text_frame.add_paragraph()
    summary_text.text = title
    summary_text.font.size = Pt(44)
    summary_text.font.bold = True
    summary_text.font.color.rgb = RGBColor(255, 255, 255)  # White font color

    # Add "Assignee/Contributors" text box with label
    assignee_left = Inches(0.5)
    assignee_top = Inches(0.9)
    assignee_width = Inches(4)
    assignee_height = Inches(0.5)
    assignee_box = slide.shapes.add_textbox(assignee_left, assignee_top, assignee_width, assignee_height)
    assignee_text_frame = assignee_box.text_frame
    assignee_text = assignee_text_frame.add_paragraph()

    # Bold "Assignee/Contributors:" label
    run1 = assignee_text.add_run()
    run1.text = "Assignee/Contributors: "
    run1.font.size = Pt(20)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0, 0, 0)  # Dark blue font color

    # Non-bold `assignee` variable
    run2 = assignee_text.add_run()
    run2.text = f"{assignee}"
    run2.font.size = Pt(20)
    run2.font.bold = False
    run2.font.color.rgb = RGBColor(0, 0, 0)  # Dark blue font color

    # Add "Stakeholder" text box with label
    stakeholder_left = Inches(7.5)
    stakeholder_top = Inches(0.9)
    stakeholder_width = Inches(4)
    stakeholder_height = Inches(0.5)
    stakeholder_box = slide.shapes.add_textbox(stakeholder_left, stakeholder_top, stakeholder_width, stakeholder_height)
    stakeholder_text_frame = stakeholder_box.text_frame
    stakeholder_text = stakeholder_text_frame.add_paragraph()

    # Bold "Stakeholder:" label
    run1 = stakeholder_text.add_run()
    run1.text = "Stakeholder: "
    run1.font.size = Pt(20)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0, 0, 0)  # Dark blue font color

    # Non-bold `stakeholder` variable
    run2 = stakeholder_text.add_run()
    run2.text = f"{stakeholder}"
    run2.font.size = Pt(20)
    run2.font.bold = False
    run2.font.color.rgb = RGBColor(0, 0, 0)  # Dark blue font color

    # Add "KRs" text box with label
    krs_left = Inches(7.5)
    krs_top = Inches(1.4)  # Adjusted slightly below the "Stakeholder" box
    krs_width = Inches(4)
    krs_height = Inches(0.5)
    krs_box = slide.shapes.add_textbox(krs_left, krs_top, krs_width, krs_height)
    krs_text_frame = krs_box.text_frame
    krs_text = krs_text_frame.add_paragraph()

    # Bold "KRs:" label
    run1 = krs_text.add_run()
    run1.text = "KRs: "
    run1.font.size = Pt(20)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0, 0, 0)  # Dark blue font color

    # Non-bold `krs` variable
    run2 = krs_text.add_run()
    run2.text = f"{krs}"
    run2.font.size = Pt(20)
    run2.font.bold = False
    run2.font.color.rgb = RGBColor(0, 0, 0)  # Dark blue font color

    # Add "Deadline" text box with label
    deadline_left = Inches(0.5)
    deadline_top = Inches(1.4)  # Adjusted slightly below the "KRs" box
    deadline_width = Inches(4)
    deadline_height = Inches(0.5)
    deadline_box = slide.shapes.add_textbox(deadline_left, deadline_top, deadline_width, deadline_height)
    deadline_text_frame = deadline_box.text_frame
    deadline_text = deadline_text_frame.add_paragraph()

    # Creating the "Deadline:" part with bold formatting
    run1 = deadline_text.add_run()
    run1.text = "Deadline: "
    run1.font.size = Pt(20)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0, 0, 0)  # Green font color

    # Creating the `deadline` part with normal (non-bold) formatting
    run2 = deadline_text.add_run()
    run2.text = f"{deadline}"
    run2.font.size = Pt(20)
    run2.font.bold = False
    run2.font.color.rgb = RGBColor(0, 0, 0)  # Green font color

    # Add a single line above both "Accomplishments" and "Challenges"
    line = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_2_ACCENT_BAR, Inches(16.25), Inches(2.4), Inches(35), 0)
    line.line.color.rgb = RGBColor(25, 94, 95)  # New color from hex #195E5F
    line.line.width = Pt(2)  # Line thickness

    # Add "Accomplishments" header
    accomplishments_title_left = Inches(0.5)
    accomplishments_title_top = Inches(2.5)
    accomplishments_title_width = Inches(4)
    accomplishments_title_height = Inches(0.5)
    accomplishments_title_box = slide.shapes.add_textbox(accomplishments_title_left, accomplishments_title_top,
                                                         accomplishments_title_width, accomplishments_title_height)
    accomplishments_title_text_frame = accomplishments_title_box.text_frame
    accomplishments_title_text_frame.text = "Accomplishments"
    accomplishments_title_text_frame.paragraphs[0].font.size = Pt(24)
    accomplishments_title_text_frame.paragraphs[0].font.bold = True
    accomplishments_title_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Green color

    # Add "Accomplishments" text box
    accomplishments_text_left = Inches(0.5)
    accomplishments_text_top = Inches(3)
    accomplishments_text_width = Inches(4)
    accomplishments_text_height = Inches(1)
    accomplishments_text_box = slide.shapes.add_textbox(accomplishments_text_left, accomplishments_text_top,
                                                        accomplishments_text_width, accomplishments_text_height)
    accomplishments_text_frame = accomplishments_text_box.text_frame
    accomplishments_text_frame.clear()  # Clear any existing content

    if accomplishments == "None":
        accomplishments = [
            "",
            "",
        ]
        # Add each accomplishment as a bullet point
        for accomplishment in accomplishments:
            p = accomplishments_text_frame.add_paragraph()
            p.text = accomplishment
            p.font.size = Pt(20)
            p.level = 0  # Bullet point level 0 (default)
            p.font.bold = False  # Optional: Set to True if you want bold text
            p.space_before = Pt(6)  # Optional: Add space before each bullet point
            p.space_after = Pt(6)  # Optional: Add space after each bullet point
            p.bullet = True  # Enable bullet points for the paragraph
    else:
        p = accomplishments_text_frame.add_paragraph()
        p.text = accomplishments
        p.font.size = Pt(20)
        p.level = 0  # Bullet point level 0 (default)
        p.font.bold = False  # Optional: Set to True if you want bold text
        p.space_before = Pt(6)  # Optional: Add space before each bullet point
        p.space_after = Pt(6)  # Optional: Add space after each bullet point



    # Add "Challenges" header
    challenges_title_left = Inches(7.5)
    challenges_title_top = Inches(2.5)
    challenges_title_width = Inches(4)
    challenges_title_height = Inches(0.5)
    challenges_title_box = slide.shapes.add_textbox(challenges_title_left, challenges_title_top,
                                                    challenges_title_width, challenges_title_height)
    challenges_title_text_frame = challenges_title_box.text_frame
    challenges_title_text_frame.text = "Challenges"
    challenges_title_text_frame.paragraphs[0].font.size = Pt(24)
    challenges_title_text_frame.paragraphs[0].font.bold = True
    challenges_title_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black color

    # Add "Challenges" text box
    challenges_text_left = Inches(7.5)
    challenges_text_top = Inches(3)
    challenges_text_width = Inches(4)
    challenges_text_height = Inches(1)
    challenges_text_box = slide.shapes.add_textbox(challenges_text_left, challenges_text_top,
                                                   challenges_text_width, challenges_text_height)
    challenges_text_frame = challenges_text_box.text_frame
    challenges_text_frame.clear()  # Clear any existing content

    if challenges == "None":
        challenges = [
            "",
            "",
        ]
        for challenge in challenges:
            p = challenges_text_frame.add_paragraph()
            p.text = challenge
            p.font.size = Pt(20)
            p.level = 0  # Bullet point level 0 (default)
            p.font.bold = False  # Optional: Set to True if you want bold text
            p.space_before = Pt(6)  # Optional: Add space before each bullet point
            p.space_after = Pt(6)  # Optional: Add space after each bullet point
            p.bullet = True  # Enable bullet points for the paragraph
    else:
        print(challenges)
        p = challenges_text_frame.add_paragraph()
        p.text = challenges
        p.font.size = Pt(20)
        p.level = 0  # Bullet point level 0 (default)
        p.font.bold = False  # Optional: Set to True if you want bold text
        p.space_before = Pt(6)  # Optional: Add space before each bullet point
        p.space_after = Pt(6)  # Optional: Add space after each bullet point

    # Add a single line above both "Accomplishments" and "Challenges"
    line = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_2_ACCENT_BAR, Inches(16.25), Inches(6), Inches(35), 0)
    line.line.color.rgb = RGBColor(25, 94, 95)  # New color from hex #195E5F
    line.line.width = Pt(2)  # Line thickness

    # Function to add a rounded rectangle around a text box
    def add_rounded_rectangle_around_textbox(slide, left, top, width, height, line_color=RGBColor(0, 0, 0)):
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        rect.line.color.rgb = line_color
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill for the rounded rectangle
        return rect

    # Function to create a centered text box
    def add_centered_textbox(slide, left, top, width, height, top_text_box, label, variable, font_size=Pt(18)):
        add_rounded_rectangle_around_textbox(slide, left, top, width, height)
        text_box = slide.shapes.add_textbox(left, top_text_box, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()  # Clear any default content

        # Center the text horizontally and vertically
        text_frame.vertical_anchor = MSO_SHAPE.RECTANGLE
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Create the bold label part
        p = text_frame.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run1 = p.add_run()
        run1.text = label
        run1.font.size = font_size
        run1.font.bold = True

        # Create the non-bold variable part
        run2 = p.add_run()
        run2.text = variable
        run2.font.size = font_size
        run2.font.bold = False

    # Add Status text box with centered label and value
    add_centered_textbox(slide, Inches(0.25), Inches(6.5), Inches(2.3), Inches(0.5), Inches(6.3), "Start: ",
                         f"{start}")

    # Add Start text box with centered label and value
    add_centered_textbox(slide, Inches(2.7), Inches(6.5), Inches(2.3), Inches(0.5), Inches(6.3), "Status: ", f"{status}")

    # Add Estimate text box with centered label and value
    add_centered_textbox(slide, Inches(5.2), Inches(6.5), Inches(2.3), Inches(0.5), Inches(6.3), "Progress: ",
                         "OnTrack")

    # Add Scope Change text box with centered label and value
    add_centered_textbox(slide, Inches(7.7), Inches(6.5), Inches(2.3), Inches(0.5), Inches(6.3), "Estimate: ",
                         f"{estimate}")

    # Add Scope Change text box with centered label and value
    add_centered_textbox(slide, Inches(10.2), Inches(6.5), Inches(2.3), Inches(0.5), Inches(6.3), "Scope Chg: ",
                         f"{scope_changed}")


    # # Function to add individual status text boxes for progress statuses with dates
    # def add_status_textbox(slide, left, top, width, height, progress_status, date):
    #     # Create a label that includes the status and date
    #     label = f"{progress_status}\n"
    #     variable = date
    #
    #     # Add the centered text box with a rounded rectangle around it
    #     add_centered_textbox(slide, left, top, width, height, Inches(6.2), label, variable, font_size=Pt(18))
    #
    # # Add individual status text boxes for progress statuses with dates
    # for i, (progress_status, date) in enumerate(progress_statuses_dates):
    #     status_box_left = Inches(7) + Inches(1.6) * i
    #     add_status_textbox(slide, status_box_left, Inches(6.5), Inches(1.5), Inches(0.7), progress_status, date)

# Example usage of the template function
# create_summary_slide(prs, title="Project Summary", assignee="John Doe", deadline="August 15", krs="Key Result 1",
#                      stakeholder="Stakeholder X",
#                      accomplishments="Completed module A", challenges="Integration issues", status="In Progress",
#                      start="July 1",
#                      estimate="5 days", scope_changed="No",
#                      progress_statuses_dates=[("On Track", "2024-08-01"), ("At Risk", "2024-08-05"),
#                                               ("Completed", "2024-08-10")])
#
# create_summary_slide(prs, title="Sprint 2 Review", assignee="Jane Doe", deadline="September 1", krs="Key Result 2",
#                      stakeholder="Stakeholder Y",
#                      accomplishments="Implemented feature B", challenges="Testing delays", status="Completed",
#                      start="August 1",
#                      estimate="7 days", scope_changed="Yes",
#                      progress_statuses_dates=[("On Track", "2024-08-01"), ("At Risk", "2024-08-05"),
#                                               ("Completed", "2024-08-10")])
#
# # Save the presentation
# prs.save('summary_template_slides_with_all_fields.pptx')


def create_squad_slide(prs, title="InStore Squad"):
    """
    Creates a slide with a colored background and centered title.

    :param prs: A PowerPoint Presentation object.
    :param title: The title text to be displayed on the slide.
    """
    # Define slide dimensions for 16:9 aspect ratio
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Add the squad slide with a colored background
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add a background color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(25, 94, 95)  # Blue background
    background.line.fill.background()  # No border

    # Add centered title text
    text_box_left = Inches(4)  # Center horizontally
    text_box_top = Inches(3)  # Center vertically
    text_box_width = Inches(5)
    text_box_height = Inches(1)
    text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
    text_frame = text_box.text_frame
    text_frame.text = title
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(44)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White font color
    paragraph.alignment = PP_ALIGN.CENTER



def create_assignee_slide(prs, assignee_name):
    """
    Creates a PowerPoint slide for a specific assignee with a colored background and centered title.

    :param prs: A PowerPoint Presentation object.
    :param assignee_name: Name of the assignee.
    """
    # Define slide dimensions for 16:9 aspect ratio
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Add the assignee slide with a colored background
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add a background color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(25, 94, 95)  # Background color
    background.line.fill.background()  # No border

    # Add centered title text
    text_box_left = Inches(4)  # Center horizontally
    text_box_top = Inches(3)  # Center vertically
    text_box_width = Inches(5)
    text_box_height = Inches(1)
    text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
    text_frame = text_box.text_frame
    text_frame.text = f"{assignee_name}"
    title_paragraph = text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White font color
    title_paragraph.alignment = PP_ALIGN.CENTER


def insert_gantt_chart_to_slide(prs, image_path):
    """
    Inserts the Gantt chart image into a new slide in the PowerPoint presentation.

    :param prs: The PowerPoint presentation object.
    :param image_path: The path of the Gantt chart image to insert.
    """
    # Add a blank slide layout
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add the Gantt chart image to the slide
    left = Inches(1)
    top = Inches(1)
    height = Inches(5.5)  # Adjust height, the width will be scaled automatically
    slide.shapes.add_picture(image_path, left, top, height=height)